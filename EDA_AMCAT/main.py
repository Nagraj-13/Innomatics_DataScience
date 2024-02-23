# from pptx import Presentation
# from pptx.util import Inches, Pt

# # Create a presentation object
# prs = Presentation()

# # Slide 1: Salary Distribution
# slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and Content layout
# title = slide.shapes.title
# content = slide.placeholders[1]

# title.text = "Salary Distribution"
# content.text = (
#     "The salary distribution appears to be right-skewed, meaning there are more employees towards the lower end of the salary range and fewer towards the higher end.\n\n"
#     "Key Statistics:\n"
#     "- Median: ₹300,000\n"
#     "- First Quartile (Q1): ₹180,000\n"
#     "- Third Quartile (Q3): ₹370,000\n"
#     "- Interquartile Range (IQR): ₹190,000\n\n"
#     "Outliers:\n"
#     "There are several outliers on the right side of the plot, with salaries exceeding ₹1 million. It's important to investigate these outliers further to understand if they are valid data points or errors."
# )

# # Slide 2: Outliers Analysis
# slide = prs.slides.add_slide(prs.slide_layouts[5])
# title = slide.shapes.title
# content = slide.placeholders[1]

# title.text = "Outliers Analysis"
# content.text = (
#     "Now it can be observed that the outliers lies in the right most part of the Q3 can be treated as rare cases, because as a fresher its highly difficult to get a salary package more than 15-20 LPA.\n"
#     "The above plots gives some valuable insights into the data. However, it is important to note that these are just preliminary results and further analysis should be conducted in order to draw more conclusions about the data."
# )

# # Slide 3: Heatmap Analysis
# slide = prs.slides.add_slide(prs.slide_layouts[5])
# title = slide.shapes.title
# content = slide.placeholders[1]

# title.text = "Heatmap Analysis"
# content.text = (
#     "So as the salary is the Target variable or object, so considering it to Bivariate Analysis. Before proceeding to bivariate analysis lets plot a heatmap to analyze the correlation matrix.\n\n"
#     "General observations:\n"
#     "- The heatmap displays the correlation coefficients between various numerical variables in your dataset, with 'Salary' as the target variable.\n"
#     "- The color scale ranges from blue (negative correlation) to yellow (positive correlation), with darker shades indicating stronger correlations.\n"
#     "- The heatmap is symmetrical, as the correlation between X and Y is the same as the correlation between Y and X."
# )

# # Slide 4: Conclusion
# slide = prs.slides.add_slide(prs.slide_layouts[5])
# title = slide.shapes.title
# content = slide.placeholders[1]

# title.text = "Conclusion"
# content.text = (
#     "Overall Trend:\n"
#     "There seems to be a weak positive correlation between salary and English, Logical, CollegeGPA, Quant score. This means that as English scores increase, salary tends to increase as well, but the relationship is not very strong.\n\n"
#     "There seems to be a positive correlation between salary and years of experience. This means that as years of experience increase, salary tends to increase as well.\n\n"
#     "The trend appears to be somewhat linear, although there is some scatter in the data points."
# )

# # Save the presentation
# prs.save("analysis_presentation.pptx")


from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Slide 1: Salary Distribution
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Salary Distribution"
content.text = (
    "- The salary distribution appears to be right-skewed, meaning there are more employees towards the lower end of the salary range and fewer towards the higher end.\n"
    "- Median: ₹300,000\n"
    "- First Quartile (Q1): ₹180,000\n"
    "- Third Quartile (Q3): ₹370,000\n"
    "- Interquartile Range (IQR): ₹190,000\n\n"
    "Outliers:\n"
    "- There are several outliers on the right side of the plot, with salaries exceeding ₹1 million. It's important to investigate these outliers further to understand if they are valid data points or errors."
)

# Slide 2: Outliers Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Outliers Analysis"
content.text = (
    "- The outliers lie in the rightmost part of Q3 and can be treated as rare cases, especially for fresher roles where it's highly unlikely to get a salary package exceeding 15-20 LPA.\n"
    "- Further analysis is required to draw more conclusions about the data."
)

# Slide 3: Heatmap Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Heatmap Analysis"
content.text = (
    "- Salary is considered the target variable for Bivariate Analysis.\n"
    "- A heatmap is used to analyze the correlation matrix between various numerical variables in the dataset, with 'Salary' as the target variable.\n"
    "- The color scale ranges from blue (negative correlation) to yellow (positive correlation), with darker shades indicating stronger correlations.\n"
    "- The heatmap is symmetrical, indicating the correlation between X and Y is the same as Y and X."
)

# Slide 4: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = (
    "- Overall Trend:\n"
    "  - There is a weak positive correlation between salary and English, Logical, CollegeGPA, and Quant scores.\n"
    "  - There is a positive correlation between salary and years of experience.\n"
    "  - The trend appears to be somewhat linear, with some scatter in the data points."
)

# Save the presentation
prs.save("analysis_presentation.pptx")
